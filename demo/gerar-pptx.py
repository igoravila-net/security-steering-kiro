#!/usr/bin/env python3
"""Gera apresentacao executiva COGNA Security Guardrails em PPTX - Tema claro com cores COGNA."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Cores COGNA
PURPLE = RGBColor(0x6B, 0x2D, 0x8B)
PURPLE_DARK = RGBColor(0x4A, 0x1D, 0x6B)
PURPLE_LIGHT = RGBColor(0xE8, 0xDE, 0xF0)
GOLD = RGBColor(0xF5, 0xA6, 0x23)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x2D, 0x2D, 0x2D)
TEXT_GRAY = RGBColor(0x66, 0x66, 0x66)
TEXT_LIGHT = RGBColor(0x99, 0x99, 0x99)
RED_SOFT = RGBColor(0xD4, 0x4B, 0x4B)
GREEN_SOFT = RGBColor(0x2D, 0x8B, 0x5E)

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

def set_slide_bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, size=18, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tf

def add_accent_bar(slide):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.06), Inches(9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PURPLE
    shape.line.fill.background()

def add_gold_line(slide, top):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(top), Inches(3), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()

# SLIDE 1: Titulo
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_accent_bar(slide)
top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(16), Inches(0.3))
top_bar.fill.solid()
top_bar.fill.fore_color.rgb = PURPLE
top_bar.line.fill.background()
add_text(slide, 1.5, 3, 13, 1.5, "Security Guardrails", size=54, bold=True, color=PURPLE_DARK)
add_gold_line(slide, 4.3)
add_text(slide, 1.5, 4.8, 13, 1, "Protecao automatizada de codigo em tempo real", size=22, color=TEXT_GRAY)
add_text(slide, 1.5, 5.8, 13, 0.5, "para o Grupo COGNA", size=18, color=TEXT_GRAY)
add_text(slide, 1.5, 7.8, 6, 0.4, "Seguranca da Informacao — Grupo COGNA", size=11, color=TEXT_LIGHT)
add_text(slide, 10, 7.8, 5, 0.4, "v2.4.2 | Maio 2026", size=11, color=TEXT_LIGHT, align=PP_ALIGN.RIGHT)

# SLIDE 2: O Problema
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_accent_bar(slide)
add_text(slide, 14, 0.6, 2, 0.4, "CONTEXTO", size=10, color=TEXT_LIGHT, align=PP_ALIGN.RIGHT)
add_text(slide, 1.5, 1, 13, 1, "Vulnerabilidades sao detectadas", size=32, color=TEXT_DARK)
add_text(slide, 1.5, 1.8, 13, 1, "tarde demais", size=32, bold=True, color=PURPLE)
add_text(slide, 1.5, 3.2, 6, 0.5, "CENARIO ATUAL", size=11, bold=True, color=RED_SOFT)
problemas = ["Encontradas no PR review ou Veracode", "Custo 10x maior a cada fase",
             "Devs nao sao especialistas em seguranca", "Dependencias com CVEs sem verificacao",
             "Politicas existem mas nao sao aplicadas"]
for i, p in enumerate(problemas):
    add_text(slide, 1.5, 3.9 + i * 0.55, 6, 0.5, f"  {p}", size=13, color=TEXT_GRAY)
add_text(slide, 8.5, 3.2, 6, 0.5, "COM SECURITY GUARDRAILS", size=11, bold=True, color=GREEN_SOFT)
solucoes = ["Bloqueadas antes de serem escritas", "Correcao automatica no momento",
            "Regras aplicadas transparentemente", "CVEs verificados antes de adicionar",
            "12 politicas COGNA automatizadas"]
for i, s in enumerate(solucoes):
    add_text(slide, 8.5, 3.9 + i * 0.55, 6, 0.5, f"  {s}", size=13, color=TEXT_GRAY)

# SLIDE 3: Numeros
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_accent_bar(slide)
add_text(slide, 14, 0.6, 2, 0.4, "IMPACTO", size=10, color=TEXT_LIGHT, align=PP_ALIGN.RIGHT)
add_text(slide, 1.5, 1, 13, 1, "Protecao em numeros", size=32, bold=True, color=TEXT_DARK)
add_gold_line(slide, 2.2)
metrics = [("46", "CWEs cobertas"), ("13", "Linguagens"), ("100%", "Codigo protegido"),
           ("26", "Hooks ativos"), ("5", "Ecossistemas SCA"), ("12", "Politicas COGNA")]
for i, (val, label) in enumerate(metrics):
    col, row = i % 3, i // 3
    x, y = 1.5 + col * 4.5, 3.2 + row * 2.6
    add_text(slide, x, y, 3.5, 1, val, size=44, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
    add_text(slide, x, y + 1, 3.5, 0.5, label, size=12, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

# SLIDE 4: Como Funciona
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_accent_bar(slide)
add_text(slide, 14, 0.6, 2, 0.4, "SOLUCAO", size=10, color=TEXT_LIGHT, align=PP_ALIGN.RIGHT)
add_text(slide, 1.5, 1, 13, 1, "Como funciona", size=32, bold=True, color=TEXT_DARK)
add_text(slide, 1.5, 2.2, 13, 0.5, "Plugin para o IDE Kiro — intercepta e corrige codigo inseguro em tempo real", size=14, color=TEXT_GRAY)
steps = [("1", "Desenvolvedor escreve codigo normalmente", "Zero friccao para codigo seguro"),
         ("2", "Power intercepta e classifica", "SKIP para docs/testes, FULL para endpoints"),
         ("3", "Se inseguro: bloqueia e corrige", "Correcao automatica com explicacao do risco"),
         ("4", "Se seguro: permite sem friccao", "89% resolvidas em menos de 1 segundo")]
for i, (num, title, desc) in enumerate(steps):
    y = 3.2 + i * 1.4
    add_text(slide, 1.5, y, 0.6, 0.6, num, size=20, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
    add_text(slide, 2.4, y, 12, 0.5, title, size=16, bold=True, color=TEXT_DARK)
    add_text(slide, 2.4, y + 0.45, 12, 0.4, desc, size=12, color=TEXT_GRAY)

# SLIDE 5: Cobertura
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_accent_bar(slide)
add_text(slide, 14, 0.6, 2, 0.4, "COBERTURA", size=10, color=TEXT_LIGHT, align=PP_ALIGN.RIGHT)
add_text(slide, 1.5, 1, 13, 1, "O que protegemos", size=32, bold=True, color=TEXT_DARK)
add_gold_line(slide, 2.2)
coverage = [("Vulnerabilidades", "46 CWEs — Top 25 MITRE 2024 + OWASP Top 10:2025"),
            ("Supply Chain", "npm, pip, Maven, NuGet, Composer — 26+ pacotes proibidos"),
            ("Infraestrutura", "Docker, Terraform, K8s — checklist de 7 itens"),
            ("LGPD", "Mascaramento de PII, consentimento, logs sem dados pessoais"),
            ("Observabilidade", "Padrao GELF COGNA, CorrelationID, niveis corretos"),
            ("CI/CD", "GitHub Actions, GitLab CI, Jenkins — secrets e permissions"),
            ("IA Segura", "LLM Top 10:2025 — prompt injection, data poisoning")]
for i, (dim, scope) in enumerate(coverage):
    y = 2.8 + i * 0.8
    add_text(slide, 1.5, y, 3, 0.5, dim, size=14, bold=True, color=PURPLE)
    add_text(slide, 4.8, y, 10, 0.5, scope, size=13, color=TEXT_GRAY)

# SLIDE 6: Exemplos
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_accent_bar(slide)
add_text(slide, 14, 0.6, 2, 0.4, "EXEMPLOS", size=10, color=TEXT_LIGHT, align=PP_ALIGN.RIGHT)
add_text(slide, 1.5, 1, 13, 1, "Protecao em acao", size=32, bold=True, color=TEXT_DARK)
add_text(slide, 1.5, 2.8, 6, 0.5, "BLOQUEADO", size=11, bold=True, color=RED_SOFT)
blocked = ["SQL Injection — concatenacao de queries", "Credenciais hardcoded no codigo",
           "Dependencias com CVEs conhecidos", "Dockerfile rodando como root",
           "Endpoints sem autenticacao", "PII em logs de producao"]
for i, b in enumerate(blocked):
    add_text(slide, 1.5, 3.5 + i * 0.55, 6, 0.5, f"  {b}", size=13, color=TEXT_GRAY)
add_text(slide, 8.5, 2.8, 6, 0.5, "CORRECAO APLICADA", size=11, bold=True, color=GREEN_SOFT)
fixed = ["Queries parametrizadas", "Vault / variaveis de ambiente", "Versoes seguras sugeridas",
         "Multi-stage + USER nao-root", "Auth + rate limiting", "Mascaramento automatico"]
for i, f in enumerate(fixed):
    add_text(slide, 8.5, 3.5 + i * 0.55, 6, 0.5, f"  {f}", size=13, color=TEXT_GRAY)

# SLIDE 7: ROI
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_accent_bar(slide)
add_text(slide, 14, 0.6, 2, 0.4, "VALOR", size=10, color=TEXT_LIGHT, align=PP_ALIGN.RIGHT)
add_text(slide, 1.5, 1, 13, 1, "Retorno sobre investimento", size=32, bold=True, color=TEXT_DARK)
add_text(slide, 2, 3, 12, 1.5, '"Corrigir em producao custa ate 100x mais\ndo que prevenir durante o desenvolvimento."', size=20, color=TEXT_GRAY)
add_gold_line(slide, 5)
roi = [("0", "Custo de licenca"), ("-83%", "Reducao de friccao"), ("0", "Configuracao necessaria")]
for i, (val, label) in enumerate(roi):
    x = 1.5 + i * 4.5
    add_text(slide, x, 5.8, 3.5, 1, val, size=44, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
    add_text(slide, x, 6.8, 3.5, 0.5, label, size=12, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

# SLIDE 8: Compliance
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_accent_bar(slide)
add_text(slide, 14, 0.6, 2, 0.4, "COMPLIANCE", size=10, color=TEXT_LIGHT, align=PP_ALIGN.RIGHT)
add_text(slide, 1.5, 1, 13, 1, "Alinhamento regulatorio", size=32, bold=True, color=TEXT_DARK)
add_gold_line(slide, 2.2)
compliance = [("OWASP Top 10:2025", "100% — todas as categorias"),
              ("API Security Top 10:2023", "100% — validacao, auth, rate limiting"),
              ("CWE Top 25 MITRE 2024", "100% — 25/25 CWEs"),
              ("LGPD (Lei 13.709/2018)", "Hook dedicado — mascaramento, consentimento"),
              ("ISO 27001:2022", "Controles tecnicos automatizados"),
              ("Politicas COGNA", "12 politicas corporativas aplicadas")]
for i, (fw, cov) in enumerate(compliance):
    y = 2.8 + i * 0.9
    add_text(slide, 1.5, y, 5, 0.5, fw, size=15, bold=True, color=TEXT_DARK)
    add_text(slide, 7, y, 8, 0.5, cov, size=13, color=TEXT_GRAY)

# SLIDE 9: Proximos Passos
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_accent_bar(slide)
add_text(slide, 14, 0.6, 2, 0.4, "ROADMAP", size=10, color=TEXT_LIGHT, align=PP_ALIGN.RIGHT)
add_text(slide, 1.5, 1, 13, 1, "Proximos passos", size=32, bold=True, color=TEXT_DARK)
add_gold_line(slide, 2.2)
roadmap = [("Piloto com squads selecionados", "Ativar em 2-3 projetos para validar"),
           ("Metricas de adocao", "Hook automatico coleta dados — dashboard de compliance"),
           ("Feedback loop Veracode", "Mapear findings para melhorar regras"),
           ("Expansao mobile", "Certificate pinning, jailbreak detection")]
for i, (title, desc) in enumerate(roadmap):
    y = 3 + i * 1.4
    bullet = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.6), Inches(y + 0.15), Inches(0.15), Inches(0.15))
    bullet.fill.solid()
    bullet.fill.fore_color.rgb = GOLD
    bullet.line.fill.background()
    add_text(slide, 2.2, y, 12, 0.5, title, size=16, bold=True, color=TEXT_DARK)
    add_text(slide, 2.2, y + 0.45, 12, 0.4, desc, size=12, color=TEXT_GRAY)

# SLIDE 10: CTA
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, PURPLE_LIGHT)
add_accent_bar(slide)
bottom_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(8.5), Inches(16), Inches(0.5))
bottom_bar.fill.solid()
bottom_bar.fill.fore_color.rgb = PURPLE
bottom_bar.line.fill.background()
add_text(slide, 1.5, 3, 13, 1.5, "Seguranca por design.", size=44, bold=True, color=PURPLE_DARK, align=PP_ALIGN.CENTER)
add_text(slide, 1.5, 4.5, 13, 1, "Transparente para o desenvolvedor.", size=28, color=TEXT_DARK, align=PP_ALIGN.CENTER)
add_gold_line(slide, 5.8)
add_text(slide, 1.5, 6.2, 13, 0.8, "100% do codigo de producao protegido, em 13 linguagens,\ncontra 46 categorias de vulnerabilidade.", size=15, color=TEXT_GRAY, align=PP_ALIGN.CENTER)
add_text(slide, 1.5, 7.5, 13, 0.5, "Seguranca da Informacao — Grupo COGNA", size=11, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

prs.save("demo/APRESENTACAO-EXECUTIVA.pptx")
print("Apresentacao gerada: demo/APRESENTACAO-EXECUTIVA.pptx")
